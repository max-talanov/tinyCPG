#!/usr/bin/env python3
"""Build a figure-review deck (12 figures in article order) as a .pptx."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image
import os

FIGDIR = "paper/figures"
OUT = "tinyCPG_figures.pptx"

SLATE = RGBColor(0x1E, 0x2A, 0x38)
TEAL  = RGBColor(0x0E, 0x7C, 0x7B)
DARK  = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ICE   = RGBColor(0xCA, 0xDC, 0xFC)

# (file, section tag, title, one-line caption)
SLIDES = [
    ("fig1_schematic", "METHODS",
     "Conceptual architecture of the two-leg spinal CPG",
     "Bilateral half-centre schematic: RG nuclei, motoneurons, muscles, Ia/cutaneous afferents, commissural interneurons."),
    ("fig_architecture_implemented", "METHODS",
     "As-implemented architecture (one leg; contralateral is the mirror)",
     "Every projection the model actually builds, colour-coded: orange excitatory, blue inhibitory, red muscle, grey inputs."),
    ("fig_connectivity", "METHODS",
     "Synaptic weight & delay distributions across all projections",
     "(a) learned-weight histograms, (b) mean±s.d. weight per projection, (c) conduction+synaptic delay per projection."),
    ("fig_stdp_weight_matrix", "RESULTS · Self-organisation",
     "STDP weights — both legs × 3 projections × 5 locomotion modes",
     "CUT→RG-E converges to ≈63 pA in every mode; the two Ia projections self-stabilise at ≈4–5 pA; L/R symmetric."),
    ("fig_stdp_weights_grid", "RESULTS · Self-organisation",
     "Weight trajectories per rehabilitation mode (5 modes × 3 λ)",
     "Learning rate sets the time-to-plateau (~7 s / ~75 s / >120 s); unloading (bottom rows) slows CUT→RG-E convergence."),
    ("fig_force_stages_lam1em3", "RESULTS · Rehabilitation progress",
     "Force at three learning stages  —  λ = 10⁻³",
     "Counter-phase force develops from ragged early bursts to a clean antiphase as the descending weight converges."),
    ("fig_force_stages_lam1em4", "RESULTS · Rehabilitation progress",
     "Force at three learning stages  —  λ = 10⁻⁴",
     "Intermediate rate: the three stages are most distinct, spanning the full 120 s convergence."),
    ("fig_force_stages_lam1em5", "RESULTS · Rehabilitation progress",
     "Force at three learning stages  —  λ = 10⁻⁵",
     "Slow rate: the descending weight barely develops, so the gait stays under-developed across all windows."),
    ("fig_network_matrix_lam1em3", "RESULTS · Network activity",
     "Full-circuit population activity  —  λ = 10⁻³",
     "All 16 populations alternate in counter-phase with a 180° interlimb offset; extensor side collapses under air stepping."),
    ("fig_network_matrix_lam1em4", "RESULTS · Network activity",
     "Full-circuit population activity  —  λ = 10⁻⁴",
     "Same organisation at the intermediate learning rate."),
    ("fig_network_matrix_lam1em5", "RESULTS · Network activity",
     "Full-circuit population activity  —  λ = 10⁻⁵",
     "At the under-converged rate the counter-phase is weaker, most visibly at higher cadence."),
    ("fig9_epidural_contrast", "RESULTS · Epidural stimulation",
     "Epidural stimulation rescues stepping under unloading",
     "Holding the paced cutaneous drive intact (stim) preserves the rhythm (corr ≈−0.98); the natural arm collapses (−0.52)."),
]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
blank = prs.slide_layouts[6]


def txt(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT,
        font="Calibri", anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.name = font
    r.font.color.rgb = color
    return tb


# --- Title slide (dark) ---
s = prs.slides.add_slide(blank)
bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = SLATE; bg.line.fill.background()
bg.shadow.inherit = False
s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
txt(s, 1.0, 2.5, 11.3, 1.5, "tinyCPG — Figure Deck", 46, WHITE, bold=True)
txt(s, 1.0, 3.7, 11.3, 1.0,
    "Self-organisation of a closed-loop spinal locomotor CPG with STDP", 22, ICE)
txt(s, 1.0, 4.5, 11.3, 1.0,
    "Figures in article order  ·  Methods + Results  ·  5 locomotion modes × 3 STDP rates",
    16, RGBColor(0x9F, 0xB3, 0xC8))

# --- Figure slides ---
BOX_X, BOX_Y, BOX_W, BOX_H = 0.45, 1.5, 12.43, 5.15
for i, (fn, tag, title, cap) in enumerate(SLIDES, 1):
    s = prs.slides.add_slide(blank)
    txt(s, 0.55, 0.33, 11.0, 0.35, tag.upper(), 13, TEAL, bold=True)
    txt(s, 0.55, 0.66, 12.2, 0.85, title, 25, DARK, bold=True)
    # fit image preserving aspect
    p = os.path.join(FIGDIR, fn + ".png")
    iw, ih = Image.open(p).size
    A = iw / ih
    if A > BOX_W / BOX_H:
        w = BOX_W; h = BOX_W / A
    else:
        h = BOX_H; w = BOX_H * A
    x = BOX_X + (BOX_W - w) / 2
    y = BOX_Y + (BOX_H - h) / 2
    s.shapes.add_picture(p, Inches(x), Inches(y), Inches(w), Inches(h))
    txt(s, 0.55, 6.95, 12.2, 0.5, cap, 12, MUTED)
    txt(s, 12.55, 7.05, 0.6, 0.3, str(i), 12, MUTED, align=PP_ALIGN.RIGHT)

prs.save(OUT)
print(f"saved {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
